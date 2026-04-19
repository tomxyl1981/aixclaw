from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(139, 69, 19)
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(28)
        p2.font.color.rgb = RGBColor(255, 223, 186)
        p2.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 250, 240)
    background.line.fill.background()
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(139, 69, 19)
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(60, 40, 20)
        p.space_after = Pt(10)

add_title_slide(prs, "四川林燊商贸有限公司", "柔性定制 · 更懂餐饮 · 合作 · 创新 · 共赢")

about = ["坐落于四川成都郫都川菜产业园","深耕餐饮调味领域十年","专注火锅底料产业链深度开发","为餐饮企业定制标准化底料","集研发、生产、定制与销售为一体","","双资质认证：","HALAL清真认证（5101240047）","出口食品生产企业备案（7900/17044）"]
add_content_slide(prs, "关于我们", about)

brand = ["核心竞争力：真材实料","","品牌承诺：","拒绝以低价竞争牺牲品质","坚持用真材实料构筑产品壁垒","食材真实、工艺真诚、数据真实","","企业愿景：成为全球餐饮调料领域的中国真味标杆","企业使命：让世界共享中国餐饮调料好味道"]
add_content_slide(prs, "旗下品牌「林燊」", brand)

why = ["深度赋能：品类调研+产品推荐+全周期跟踪","","全品类覆盖：","火锅调味料全品类生产","多味型适配区域市场","欧盟级原料溯源+零添加工艺","千家门店实证复购率91%","","真诚利他：","供应链资源对接（不拿回扣）","免费一对一技术指导服务"]
add_content_slide(prs, "为什么选择我们？", why)

series = [
    ("01 | 牛油火锅系列", ["牛油9号/A号：新派畅销款，性价比高","牛油老火锅2/5/7号：重庆风味，重油重辣","重庆水火锅：可以喝汤的火锅","藤椒牛油底料：清一色藤椒风味","真香锅无渣牛油油料","浓香/酱香/兼香型牛油：混合一体化包装","经典牛油混合料M号：三合一包装"]),
    ("02 | 清油火锅系列", ["2号清油底料：清爽不上火","4号清油底料：纯菜籽油，适合东三省","经典清油混合料W号：三合一包装","真香锅清油火锅底料：适合地摊、中餐"]),
    ("03 | 三合一系列（清汤/酸汤）", ["番茄底料2号/黄金板栗鸡底料","2号松茸菌汤/骨汤膏/黄金松茸牛肝菌","贵州红酸汤/泰式冬阴功","酸辣金汤/经典酸菜/酸汤肥牛","牛骨汤/翘脚牛骨汤/草本汤锅"]),
    ("04 | 鱼蛙火锅系列", ["美蛙鱼底料：爆款红汤麻辣，高性价比","麻辣香水鱼底料/霸王鱼底料","青椒椒麻底料/金汤泡椒鱼底料","酸萝卜鱼底料/酸菜鱼底料","贵州酸汤鱼底料"]),
    ("05 | 鸡兔牛杂系列", ["渣渣牛肉火锅底料","混合油牛杂火锅底料","鸡兔牛油底料/鸡兔混合油底料","猪肚鸡底料/花胶鸡底料","柴火鸡底料/鸡公煲底料"]),
    ("06 | 串串火锅系列", ["混合油热锅串串底料：30%牛油+70%清油","纯牛油热锅串串底料","混合油冷锅串串底料","纯清油冷锅串串底料"]),
    ("07 | 麻辣烫系列", ["麻辣烫底料/盘盘麻辣烫底料","番茄麻辣烫底料/藤椒麻辣烫底料","酸辣金汤麻辣烫底料/三鲜骨汤麻辣烫底料","北方麻辣烫底料：含奶粉配方"]),
    ("08 | 冒菜系列", ["冒菜底料/冒烤鸭底料/干拌冒菜底料","番茄/藤椒/酸辣金汤/三鲜骨汤冒菜","贵州酸汤冒菜/冬阴功冒菜/酱香卤味冒菜"]),
    ("09 | 钵钵鸡系列", ["钵钵鸡红油：渣少，乐山特色","藤椒钵钵鸡底料/酸辣金汤钵钵鸡底料","冷沾沾红油：辣椒香，无渣"]),
    ("10 | 火锅蘸酱系列", ["菌王酱、豆豉酱、香辣酱、香辣牛肉酱","沙茶酱、XO酱、青椒椒麻酱、蒜蓉酱","适用于火锅、烧烤、串串等"]),
    ("11 | 小龙虾系列", ["香辣小龙虾底料/蒜泥小龙虾底料","油焖小龙虾底料/十三香小龙虾底料","辣卤小龙虾底料"]),
    ("12 | 辣卤系列", ["辣卤底料：经典辣卤、香卤","应用于现捞卤味"]),
    ("13 | 干锅系列", ["香辣干锅底料：适用虾蟹、鱼蛙、鸡兔等","酱香干锅底料/藤椒干锅底料","香辣孜然干锅底料"]),
    ("14 | 烤鱼系列", ["香辣/香辣孜然/青椒椒麻烤鱼","酸汤金汤/蒜泥/酱香烤鱼","豆豉/泡椒/番茄/贵州红酸汤烤鱼","酸萝卜/冬阴功/经典酸菜/藿香烤鱼"]),
    ("15 | 米线/粉面系列", ["麻辣/酸辣金汤/三鲜清汤米线","藤椒/番茄/经典酸菜米线","泡椒/贵州红酸汤米线底料","酸辣粉红油"]),
    ("16 | 砂锅/水煮系列", ["麻辣砂锅/酸辣金汤砂锅/三鲜清汤砂锅","青花椒水煮/金汤水煮/水煮肉片"]),
    ("17 | 中餐炒菜/烧菜系列", ["凉拌红油/小炒汁/爆炒酱","麻辣香锅底料/烧菜底料/回锅肉底料","麻婆豆腐底料/毛血旺底料/蹄花骨汤料"]),
    ("18 | 通用辅料系列", ["蘸料撒料：辣椒干碟面、麻辣烧烤撒料、孜然五香烧烤撒料","腌料粉类：调味粉、腌肉粉、腌鱼粉","油类：香油、腌肉红油"])
]

for t, items in series:
    add_content_slide(prs, t, items)

strength = ["累计服务超 10000+ 餐饮门店","","覆盖业态：","火锅 | 鱼火锅 | 冒菜 | 串串 | 烤鱼 | 小龙虾 | 中餐 | 粉面小吃","","核心优势：","10年专注 | 1万+门店实证 | 自有工厂 | 资质齐全 | 柔性定制"]
add_content_slide(prs, "实力见证", strength)

contact = ["公司名称：四川林燊商贸有限公司","工厂地址：四川省成都市郫都区永盛村300号","联系人：张女士","联系电话：13438486575","","欢迎致电咨询，获取免费样品及一对一技术指导服务"]
add_content_slide(prs, "联系我们", contact)

prs.save('/home/xiaoyao/.openclaw/workspace/林燊商贸产品宣传册.pptx')
print("PPT已生成！共", len(prs.slides), "页")
